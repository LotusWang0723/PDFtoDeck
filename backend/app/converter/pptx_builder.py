"""PPT builder: generate .pptx from extracted PDF elements."""

import fitz  # for rendering curved vectors as PNG
from PIL import Image
import numpy as np
import io
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO

from .models import PageElements, ElementType, VectorElement, PathNode


# PDF points to EMU (English Metric Units, used by python-pptx)
# 1 point = 12700 EMU
PT_TO_EMU = 12700


def _pt_to_emu(val: float) -> int:
    return int(val * PT_TO_EMU)


def _rgb(color_tuple: tuple) -> RGBColor:
    r, g, b = color_tuple[:3]
    return RGBColor(min(r, 255), min(g, 255), min(b, 255))


def _add_text_element(slide, te, page_height: float):
    """Add a text box to the slide.

    pymupdf uses top-left origin (same as PPT), so NO Y-axis flip needed.
    """
    left = _pt_to_emu(te.bbox.x0)
    top = _pt_to_emu(te.bbox.y0)
    width = _pt_to_emu(te.bbox.width)
    height = _pt_to_emu(te.bbox.height)

    # Ensure minimum dimensions
    width = max(width, _pt_to_emu(10))
    height = max(height, _pt_to_emu(te.font_size * 1.5))

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = False
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = te.text
    p.font.size = Pt(te.font_size)
    p.font.bold = te.bold
    p.font.italic = te.italic
    p.font.color.rgb = _rgb(te.color)
    if te.font_name:
        p.font.name = te.font_name


def _add_image_element(slide, ie, page_height: float,
                       pdf_path: str = "", page_num: int = 0,
                       page_width: float = 0, page_height_val: float = 0):
    """Add an image to the slide.
    
    For small icon-sized images, clips from original PDF with background
    removal for cleaner rendering.
    """
    left = _pt_to_emu(ie.bbox.x0)
    top = _pt_to_emu(ie.bbox.y0)
    width = _pt_to_emu(ie.bbox.width)
    height = _pt_to_emu(ie.bbox.height)

    # Check if this is a small icon that would benefit from clip+bg-removal
    page_area = page_width * page_height_val if page_width and page_height_val else 0
    img_area = ie.bbox.width * ie.bbox.height
    is_small_icon = page_area > 0 and (img_area / page_area) < 0.05

    if is_small_icon and pdf_path:
        png = _clip_from_pdf(pdf_path, page_num, ie.bbox)
        if png:
            slide.shapes.add_picture(BytesIO(png), left, top, width, height)
            return

    # Default: use original image bytes
    stream = BytesIO(ie.image_bytes)
    slide.shapes.add_picture(stream, left, top, width, height)


def _add_freeform_shape(slide, ve: VectorElement, page_height: float):
    """Add an editable freeform shape (Plan A) to the slide.

    Uses python-pptx 1.x API: build_freeform → move_to / add_line_segments.
    pymupdf uses top-left origin — NO Y-axis flip.
    """
    if not ve.nodes:
        return

    left = _pt_to_emu(ve.bbox.x0)
    top = _pt_to_emu(ve.bbox.y0)
    width = _pt_to_emu(ve.bbox.width) or _pt_to_emu(1)
    height = _pt_to_emu(ve.bbox.height) or _pt_to_emu(1)

    bw = ve.bbox.width or 1
    bh = ve.bbox.height or 1

    def _normalize(node):
        """Normalize node coords to 0-1 range within bbox."""
        nx = (node.x - ve.bbox.x0) / bw
        ny = (node.y - ve.bbox.y0) / bh
        return (nx, ny)

    first = ve.nodes[0]
    sx, sy = _normalize(first)

    builder = slide.shapes.build_freeform(sx, sy, scale=(width, height))

    segment: list[tuple[float, float]] = []

    for node in ve.nodes[1:]:
        nx, ny = _normalize(node)

        if node.kind == "move":
            if segment:
                builder.add_line_segments(segment, close=False)
                segment = []
            builder.move_to(nx, ny)
        else:
            segment.append((nx, ny))

    if segment:
        builder.add_line_segments(segment, close=True)

    shape = builder.convert_to_shape(left, top)

    # Apply colors
    if ve.fill_color and ve.fill_color != (0, 0, 0):
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(ve.fill_color)
    else:
        shape.fill.background()

    if ve.stroke_color:
        shape.line.color.rgb = _rgb(ve.stroke_color)
        shape.line.width = Pt(ve.stroke_width or 1.0)


def _add_vector_as_image(slide, ve: VectorElement, page_height: float):
    """Add vector as a rectangle shape (Plan B fallback)."""
    left = _pt_to_emu(ve.bbox.x0)
    top = _pt_to_emu(ve.bbox.y0)
    width = _pt_to_emu(ve.bbox.width) or _pt_to_emu(10)
    height = _pt_to_emu(ve.bbox.height) or _pt_to_emu(10)

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height,
    )
    if ve.fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(ve.fill_color)
    if ve.stroke_color:
        shape.line.color.rgb = _rgb(ve.stroke_color)


def _render_curved_vector_as_png(ve: VectorElement, scale: float = 3.0) -> bytes | None:
    """Render a curved vector by clipping from original PDF page + removing background.
    
    This is a fallback that gets called with the PDF doc and page context.
    See _clip_from_pdf() for the actual implementation.
    Returns None — the real work is done via _clip_from_pdf in build_pptx.
    """
    return None


def _clip_from_pdf(
    pdf_path: str, page_num: int, bbox,
    text_bboxes: list = None,
    scale: float = 4.0,
) -> bytes | None:
    """Clip a region from the original PDF page and remove background color.
    
    Uses border pixel statistics for adaptive background detection,
    then Euclidean color distance for precise removal.
    
    Args:
        pdf_path: Path to the PDF file.
        page_num: Page index.
        bbox: BBox object with x0, y0, x1, y1.
        text_bboxes: List of BBox for text spans to mask out.
        scale: Render scale (higher = sharper).
    
    Returns:
        PNG bytes with transparent background, or None on failure.
    """
    try:
        doc = fitz.open(pdf_path)
        page = doc[page_num]
        
        clip = fitz.Rect(bbox.x0 - 1, bbox.y0 - 1, bbox.x1 + 1, bbox.y1 + 1)
        mat = fitz.Matrix(scale, scale)
        pix = page.get_pixmap(matrix=mat, clip=clip, alpha=False)
        doc.close()
        
        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        arr = np.array(img).astype(float)
        h, w = arr.shape[:2]
        
        # Mask out text regions before background detection
        if text_bboxes:
            for tb in text_bboxes:
                px0 = max(0, int((tb.x0 - bbox.x0 + 1) * scale) - 2)
                py0 = max(0, int((tb.y0 - bbox.y0 + 1) * scale) - 2)
                px1 = min(w, int((tb.x1 - bbox.x0 + 1) * scale) + 2)
                py1 = min(h, int((tb.y1 - bbox.y0 + 1) * scale) + 2)
                # Fill text areas with border color (will be computed below)
                # For now mark with NaN, handle after bg detection
        
        # Sample border pixels (outermost 3 rows/cols) for background color
        border_width = min(3, h // 4, w // 4)
        top_border = arr[:border_width, :].reshape(-1, 3)
        bottom_border = arr[-border_width:, :].reshape(-1, 3)
        left_border = arr[:, :border_width].reshape(-1, 3)
        right_border = arr[:, -border_width:].reshape(-1, 3)
        
        # Use the border side with lowest variance (most uniform = most likely pure bg)
        borders = [top_border, bottom_border, left_border, right_border]
        best_border = min(borders, key=lambda b: np.std(b))
        
        bg_mean = best_border.mean(axis=0)
        bg_std = best_border.std(axis=0)
        
        # Conservative adaptive tolerance: cap at 60 to avoid eating foreground
        tolerance = min(60, max(35, int(np.max(bg_std) * 3.5)))
        
        # Convert to RGBA
        img_rgba = img.convert("RGBA")
        arr_rgba = np.array(img_rgba)
        
        # Paint over text regions with background color
        if text_bboxes:
            bg_rgb = bg_mean.astype(np.uint8)
            for tb in text_bboxes:
                px0 = max(0, int((tb.x0 - bbox.x0 + 1) * scale) - 4)
                py0 = max(0, int((tb.y0 - bbox.y0 + 1) * scale) - 4)
                px1 = min(w, int((tb.x1 - bbox.x0 + 1) * scale) + 4)
                py1 = min(h, int((tb.y1 - bbox.y0 + 1) * scale) + 4)
                arr_rgba[py0:py1, px0:px1, :3] = bg_rgb
                arr[py0:py1, px0:px1, :] = bg_mean  # Update float array too
        
        # Euclidean color distance from background mean
        diff = np.sqrt(np.sum((arr - bg_mean.reshape(1, 1, 3)) ** 2, axis=2))
        
        # Make background transparent
        arr_rgba[diff < tolerance, 3] = 0
        
        result = Image.fromarray(arr_rgba)
        buf = io.BytesIO()
        result.save(buf, format="PNG")
        return buf.getvalue()
    except Exception:
        return None


def _add_curved_vector_as_png(slide, ve: VectorElement, page_height: float,
                               pdf_path: str = "", page_num: int = 0,
                               text_elements: list = None):
    """Clip curved vector from original PDF as high-fidelity transparent PNG.
    
    Text spans that overlap the vector area are masked out from the PNG,
    since they'll be added as separate editable text boxes.
    """
    png = None
    if pdf_path:
        # Find text bboxes that overlap this vector
        text_bboxes = []
        if text_elements:
            for te in text_elements:
                if (te.bbox.x0 < ve.bbox.x1 and te.bbox.x1 > ve.bbox.x0 and
                    te.bbox.y0 < ve.bbox.y1 and te.bbox.y1 > ve.bbox.y0):
                    text_bboxes.append(te.bbox)
        
        png = _clip_from_pdf(pdf_path, page_num, ve.bbox,
                            text_bboxes=text_bboxes)
    
    if png:
        left = _pt_to_emu(ve.bbox.x0)
        top = _pt_to_emu(ve.bbox.y0)
        width = _pt_to_emu(ve.bbox.width) or _pt_to_emu(10)
        height = _pt_to_emu(ve.bbox.height) or _pt_to_emu(10)
        slide.shapes.add_picture(BytesIO(png), left, top, width, height)
    else:
        _add_freeform_shape(slide, ve, page_height)


def _add_branding_footer(slide, page_width: float, page_height: float):
    """Add a subtle PDFtoDeck branding footer to the slide."""
    text = "Converted by PDFtoDeck"
    font_size = 6
    box_width = 120
    box_height = 12

    left = _pt_to_emu(page_width - box_width - 8)
    top = _pt_to_emu(page_height - box_height - 4)
    width = _pt_to_emu(box_width)
    height = _pt_to_emu(box_height)

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = False
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(160, 160, 180)
    p.font.italic = True


def build_pptx(
    pages: list[PageElements],
    output_path: str,
    source_filename: str = "",
    pdf_path: str = "",
) -> str:
    """Build a .pptx file from parsed PDF pages.

    Args:
        pages: List of PageElements from pdf_parser.
        output_path: Where to save the .pptx file.
        source_filename: Original PDF filename (for PPTX metadata).
        pdf_path: Path to original PDF (for high-fidelity clip rendering).

    Returns:
        The output file path.
    """
    prs = Presentation()

    # Set document metadata
    prs.core_properties.title = source_filename.replace(".pdf", "") if source_filename else "PDFtoDeck"
    prs.core_properties.comments = "Converted by PDFtoDeck — https://pdftodeck.com"

    for page in pages:
        # Set slide dimensions to match PDF page
        prs.slide_width = _pt_to_emu(page.width)
        prs.slide_height = _pt_to_emu(page.height)

        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # === Z-ORDER MATTERS ===
        # 1. Large vectors first (backgrounds, decorative shapes)
        for ve in page.vectors:
            if ve.element_type == ElementType.VECTOR_LARGE:
                _add_vector_as_image(slide, ve, page.height)

        # 2. Icon-image fallbacks (medium complexity vectors)
        for ve in page.vectors:
            if ve.element_type == ElementType.ICON_IMAGE:
                _add_vector_as_image(slide, ve, page.height)

        # 3. Editable freeform icon shapes (or PNG for curved ones)
        for ve in page.vectors:
            if ve.element_type == ElementType.ICON_SHAPE:
                if ve.has_curves:
                    _add_curved_vector_as_png(
                        slide, ve, page.height,
                        pdf_path=pdf_path, page_num=page.page_num,
                        text_elements=page.texts,
                    )
                else:
                    _add_freeform_shape(slide, ve, page.height)

        # 4. Images (on top of vector backgrounds)
        for ie in page.images:
            _add_image_element(slide, ie, page.height,
                              pdf_path=pdf_path, page_num=page.page_num,
                              page_width=page.width, page_height_val=page.height)

        # 5. Text elements (topmost, always readable)
        for te in page.texts:
            _add_text_element(slide, te, page.height)

        # 6. PDFtoDeck branding (very top)
        _add_branding_footer(slide, page.width, page.height)

    prs.save(output_path)
    return output_path
